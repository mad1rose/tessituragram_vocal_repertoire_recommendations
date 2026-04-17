"""
Build the CONDENSED 8-slide CADSCOM 2026 presentation (.pptx).

Run from the repo root:
    python CADSCOM_PRESENTATION/build_presentation_condensed.py

Requires: python-pptx, Pillow
"""

from __future__ import annotations

import io
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = pathlib.Path(__file__).resolve().parent.parent
OUT = REPO / "CADSCOM_PRESENTATION" / "CADSCOM_2026_Presentation_condensed.pptx"

# ---------------------------------------------------------------------------
# Palette (identical to full deck)
# ---------------------------------------------------------------------------
BURGUNDY   = RGBColor(0x0B, 0x1F, 0x3B)
CREAM      = RGBColor(0xFF, 0xFF, 0xFF)
GOLD       = RGBColor(0xFF, 0x5D, 0x7A)
CHARCOAL   = RGBColor(0x1A, 0x1A, 0x1A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GOLD = RGBColor(0xF4, 0xF6, 0xF8)
SOFT_GRAY  = RGBColor(0x2E, 0x4A, 0x7D)
TEAL       = RGBColor(0x3F, 0xA7, 0xA3)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers (identical to full deck)
# ---------------------------------------------------------------------------

def _set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_burgundy_header(slide, height=Inches(1.15)):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = BURGUNDY
    shp.line.fill.background()
    return shp


def _add_gold_accent_line(slide, top, width=Inches(2.5)):
    left = Inches(0.75)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = GOLD
    shp.line.fill.background()


def _add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _set_font(run, size=Pt(22), bold=False, italic=False,
              color=CHARCOAL, name="Calibri"):
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = name


def _add_title_in_header(slide, text, top=Inches(0.25)):
    tb = _add_textbox(slide, Inches(0.75), top, Inches(11.5), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, Pt(32), bold=True, color=WHITE)
    p.alignment = PP_ALIGN.LEFT
    return tb


def _add_notes(slide, text: str):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def _add_image_right(slide, img_path: str, left, top, max_w, max_h):
    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    w = max_w
    h = int(w / aspect)
    if h > max_h:
        h = max_h
        w = int(h * aspect)
    slide.shapes.add_picture(img_path, left, top, w, h)


def _add_image_tau_panel_from_alpha_figure(
    slide, img_path: str, left, top, max_w, max_h, x_start_frac: float = 0.505,
):
    """Right subplot only from alpha_sensitivity_hr1_mrr_tau.png (τ vs α)."""
    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        iw, ih = im.size
        x0 = min(max(0, int(iw * x_start_frac)), iw - 2)
        cropped = im.crop((x0, 0, iw, ih))
    ciw, cih = cropped.size
    aspect = ciw / cih
    w = max_w
    h = int(w / aspect)
    if h > max_h:
        h = max_h
        w = int(h * aspect)
    bio = io.BytesIO()
    cropped.save(bio, format="PNG")
    bio.seek(0)
    slide.shapes.add_picture(bio, left, top, w, h)


def _add_image_centered(slide, img_path: str, top=Inches(1.6),
                        max_w=Inches(11.0), max_h=Inches(5.3)):
    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    w = max_w
    h = int(w / aspect)
    if h > max_h:
        h = max_h
        w = int(h * aspect)
    left = int((SLIDE_W - w) / 2)
    slide.shapes.add_picture(img_path, left, top, w, h)


def _format_table_cell(cell, text, size, bold=False, color=CHARCOAL,
                       align=PP_ALIGN.CENTER, bg=None):
    cell.text = text
    tf = cell.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for p in tf.paragraphs:
        p.alignment = align
        for run in p.runs:
            _set_font(run, size, bold=bold, color=color)
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def _two_col_card(slide, left_x, top, w, h, border_rgb, fill_rgb,
                  title, title_color, items, item_color, item_size=Pt(17)):
    """Reusable two-column card with a bold title and bulleted items."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_x, top, w, h,
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_rgb
    if fill_rgb == BURGUNDY:
        box.line.fill.background()
    else:
        box.line.color.rgb = border_rgb
        box.line.width = Pt(2.0)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.15)

    ph = tf.paragraphs[0]
    ph.alignment = PP_ALIGN.LEFT
    rh = ph.add_run()
    rh.text = title
    _set_font(rh, Pt(20), bold=True, color=title_color)

    for item in items:
        p = tf.add_paragraph()
        p.space_before = Pt(12)
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "  \u2022  " + item
        _set_font(r, item_size, color=item_color)


# ---------------------------------------------------------------------------
# Slide 1 -- Title (identical to full deck)
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BURGUNDY)

    for y in (Inches(1.8), Inches(5.6)):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(2), y, Inches(9.333), Inches(0.04),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = GOLD
        shp.line.fill.background()

    tb = _add_textbox(slide, Inches(1), Inches(2.1), Inches(11.333), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = (
        "Content-Based Vocal Repertoire Ranking Framework\n"
        "Using Duration-Weighted Pitch Distributions"
    )
    _set_font(run, Pt(34), bold=True, color=WHITE)

    tb2 = _add_textbox(slide, Inches(1), Inches(4.4), Inches(11.333), Inches(1.0))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Madeline Johnson, Flint Million, Rajeev Bukralia"
    _set_font(run2, Pt(24), color=LIGHT_GOLD)

    tb3 = _add_textbox(slide, Inches(1), Inches(5.05), Inches(11.333), Inches(0.5))
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Minnesota State University, Mankato"
    _set_font(run3, Pt(20), italic=True, color=LIGHT_GOLD)

    tb4 = _add_textbox(slide, Inches(1), Inches(5.9), Inches(11.333), Inches(0.5))
    p4 = tb4.text_frame.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run()
    run4.text = "CADSCOM 2026"
    _set_font(run4, Pt(20), bold=True, color=GOLD)

    _add_notes(slide, (
        "[~20 seconds]\n\n"
        "Good morning/afternoon, everyone. My name is Madeline Johnson, "
        "and I am currently pursuing a master's in data science here at "
        "Minnesota State University, Mankato. My undergraduate background "
        "is in vocal performance -- I hold a Bachelor of Music from "
        "Illinois Wesleyan University -- so this project sits right at "
        "the intersection of my two worlds. This work is co-authored with "
        "Flint Million and Rajeev Bukralia, and today I want to show you "
        "how we can use data from musical scores to help singers find "
        "repertoire that actually fits their voice."
    ))


# ---------------------------------------------------------------------------
# Slide 2 -- Problem & Motivation  (merges old 2 + 3)
# ---------------------------------------------------------------------------

def slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "Problem & Motivation")
    _add_gold_accent_line(slide, Inches(1.18))

    gap = Inches(0.35)
    card_w = (Inches(12.2) - gap) / 2
    left_x = Inches(0.55)
    right_x = left_x + card_w + gap
    card_top = Inches(1.38)
    card_h = Inches(4.45)

    _two_col_card(
        slide, left_x, card_top, card_w, card_h,
        TEAL, LIGHT_GOLD,
        "The wrong piece is an injury risk",
        BURGUNDY,
        [
            "Singing outside your comfort zone strains the voice.",
            "Singers rely on intuition and trial and error to choose repertoire.",
            "Data could make this safer and more objective.",
        ],
        CHARCOAL,
    )

    _two_col_card(
        slide, right_x, card_top, card_w, card_h,
        GOLD, BURGUNDY,
        "Current filters fall short",
        GOLD,
        [
            "Range: shows extremes, not where the voice actually lives.",
            "Fach: labels inconsistent across regions (Schloneger et al., 2024).",
            "Neither captures where the voice spends its time.",
        ],
        WHITE,
    )

    tb = _add_textbox(slide, Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Neither captures where the voice actually spends its time."
    _set_font(r, Pt(20), italic=True, color=BURGUNDY)

    fn = _add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4))
    fn_p = fn.text_frame.paragraphs[0]
    fn_p.alignment = PP_ALIGN.RIGHT
    fn_run = fn_p.add_run()
    fn_run.text = "Apfelbach, 2022; Phyland et al., 1999"
    _set_font(fn_run, Pt(12), italic=True, color=SOFT_GRAY)

    _add_notes(slide, (
        "[~90 seconds]\n\n"
        "As a trained singer, one of the first things I learned is that "
        "choosing the wrong piece is not just an inconvenience -- it can "
        "lead to real vocal injury. For example, a singer with a naturally "
        "lower voice will usually not feel as comfortable spending long "
        "periods in a high range as a singer with a naturally higher voice. "
        "Research shows that when a piece's demands do not match a singer's "
        "capabilities, the risk of strain goes up significantly (Apfelbach, "
        "2022; Phyland et al., 1999). In practice, singers and voice "
        "teachers rely on intuition, trial and error, or vocal "
        "classification systems that are not consistently defined. That is "
        "what motivated this research: using data to make the process of "
        "choosing repertoire more objective and safer.\n\n"
        "Today, singers mainly rely on two approaches. The first is "
        "filtering by range -- you set a minimum and maximum pitch. That "
        "is useful, but it only tells you the extremes. A piece with one "
        "very high note at the climax looks the same as a piece that sits "
        "up high the entire time. The second approach is Fach -- the "
        "system opera uses to categorize voices by range, weight, and "
        "color. But Fach labels are not consistently defined across "
        "regions or pedagogical traditions (Schloneger et al., 2024), and "
        "many voices fall between categories. Neither approach captures "
        "where the voice actually spends its time -- and that distinction "
        "matters."
    ))


# ---------------------------------------------------------------------------
# Slide 3 -- Tessituragrams & the Gap  (merges old 4 + 5)
# ---------------------------------------------------------------------------

def slide_03_idea(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "Tessituragrams & the Gap")
    _add_gold_accent_line(slide, Inches(1.18))

    gap = Inches(0.35)
    card_w = (Inches(12.2) - gap) / 2
    left_x = Inches(0.55)
    right_x = left_x + card_w + gap
    card_top = Inches(1.38)
    card_h = Inches(4.45)

    _two_col_card(
        slide, left_x, card_top, card_w, card_h,
        SOFT_GRAY, LIGHT_GOLD,
        "Tessituragram",
        BURGUNDY,
        [
            "Thurmer (1988): histogram of singing time per pitch \u2014 a song\u2019s fingerprint.",
            "Titze (2008): duration-weighted \u2014 longer notes count more.",
            "Captures what range and Fach miss: the internal pitch-duration distribution.",
        ],
        CHARCOAL,
    )

    _two_col_card(
        slide, right_x, card_top, card_w, card_h,
        GOLD, BURGUNDY,
        "The Gap",
        GOLD,
        [
            "Existing tools (Tessa, Kassia) are purely analytic.",
            "No query-side recommendation: \u201cFind songs on THESE notes.\u201d",
            "217,000+ cataloged settings \u2014 manual search can\u2019t scale.",
        ],
        WHITE,
    )

    tb = _add_textbox(slide, Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Tessituragrams capture vocal demand \u2014 but no tool uses them to recommend."
    _set_font(r, Pt(20), italic=True, bold=True, color=BURGUNDY)

    _add_notes(slide, (
        "[~90 seconds]\n\n"
        "There is a concept that does capture where the voice spends its "
        "time: tessitura. Range tells you the highest and lowest notes, "
        "but tessitura tells you where the voice actually lives -- the "
        "pitches where it spends most of its time. A piece might touch "
        "its highest note once at the climax, but if it sits in the "
        "middle of the range for most of its duration, that is where the "
        "real demand is.\n\n"
        "A researcher named Stefan Thurmer formalized a way to represent "
        "tessitura visually. In 1988 he introduced the tessituragram -- "
        "a histogram of singing time per pitch. Titze, in 2008, refined "
        "this with duration weighting, meaning longer notes count more "
        "because sustaining a pitch is more demanding than a passing "
        "tone. The result is a fingerprint of where the voice spends its "
        "time in that piece.\n\n"
        "There are tools that work with tessituragrams, but they are "
        "purely analytic. Tessa extracts a tessituragram from a digital "
        "score. The Kassia Database displays a human-assessed tessitura "
        "for each entry. Both help you evaluate a piece after you have "
        "already found it. But neither one lets you use a tessituragram "
        "as a query -- to search for new pieces whose pitch distribution "
        "is likely to fit. The LiederNet Archive alone catalogs over "
        "217,000 art-song settings. To our knowledge, no existing system "
        "uses tessituragrams on the query side to recommend repertoire. "
        "That is the gap we wanted to explore."
    ))


# ---------------------------------------------------------------------------
# Slide 4 -- Methodology  (merges old 6 + 7)
# ---------------------------------------------------------------------------

def slide_04_method(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "Methodology")
    _add_gold_accent_line(slide, Inches(1.18))

    # Data-source banner
    src_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(1.3), Inches(11.8), Inches(0.45),
    )
    src_box.fill.solid()
    src_box.fill.fore_color.rgb = LIGHT_GOLD
    src_box.line.color.rgb = TEAL
    src_box.line.width = Pt(1.2)
    tf_src = src_box.text_frame
    tf_src.word_wrap = True
    tf_src.margin_left = Inches(0.2)
    tf_src.margin_top = Inches(0.06)
    p_src = tf_src.paragraphs[0]
    p_src.alignment = PP_ALIGN.CENTER
    r_src = p_src.add_run()
    r_src.text = "MusicXML parsed with music21: pitch \u2192 duration in quarter-note beats."
    _set_font(r_src, Pt(14), italic=True, color=CHARCOAL)

    # Pipeline strip
    steps = [
        ("Singer\nPreferences", "Range, favorites, and avoids"),
        ("Range\nFilter", "Drop songs outside range"),
        ("Ideal Profile\nConstruction", "Target profile from preferences"),
        ("Cosine Similarity\nScoring", "Similarity minus avoid penalty"),
        ("Ranked\nSong List", "Best match to worst"),
    ]

    box_w = Inches(1.85)
    box_h = Inches(1.55)
    gap_x = Inches(0.38)
    total_w = len(steps) * box_w + (len(steps) - 1) * gap_x
    start_left = int((SLIDE_W - total_w) / 2)
    top_box = Inches(1.92)
    top_desc = Inches(3.58)

    for i, (label, desc) in enumerate(steps):
        left = start_left + i * (box_w + gap_x)

        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top_box, box_w, box_h,
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = BURGUNDY
        shp.line.fill.background()
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        _set_font(run, Pt(14), bold=True, color=WHITE)
        shp.text_frame.paragraphs[0].space_before = Pt(14)

        tb = _add_textbox(slide, left, top_desc, box_w, Inches(1.0))
        tf2 = tb.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        _set_font(run2, Pt(12), color=CHARCOAL)

        if i < len(steps) - 1:
            arrow_left = left + box_w
            arrow_top = top_box + box_h / 2 - Inches(0.12)
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top,
                gap_x, Inches(0.24),
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = TEAL
            arr.line.fill.background()

    # Formula banner
    formula_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(4.72), Inches(11.8), Inches(0.72),
    )
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = BURGUNDY
    formula_box.line.fill.background()
    tf_f = formula_box.text_frame
    tf_f.word_wrap = True
    tf_f.margin_top = Inches(0.1)
    pf = tf_f.paragraphs[0]
    pf.alignment = PP_ALIGN.CENTER
    rf = pf.add_run()
    rf.text = "final_score  =  cosine_similarity(song, ideal)  \u2212  \u03B1  \u00D7  avoid_penalty"
    _set_font(rf, Pt(20), bold=True, color=WHITE, name="Consolas")

    # Three inline component labels
    lbl_top = Inches(5.65)
    lbl_h = Inches(1.4)
    lbl_w = Inches(3.7)
    lbl_gap = Inches(0.35)
    total_lbl = 3 * lbl_w + 2 * lbl_gap
    lbl_start = int((SLIDE_W - total_lbl) / 2)

    labels = [
        ("Cosine Similarity", TEAL, "Same pitches, same proportions? 1.0 = perfect match."),
        ("Avoid Penalty", GOLD, "Time on pitches to avoid. Higher = worse fit."),
        ("\u03B1 = 0.5", SOFT_GRAY, "Balances good-fit reward vs. avoid penalty."),
    ]
    for i, (title, accent, desc) in enumerate(labels):
        left = lbl_start + i * (lbl_w + lbl_gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, lbl_top, lbl_w, lbl_h,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GOLD
        card.line.color.rgb = accent
        card.line.width = Pt(2.0)
        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.15)
        tf_c.margin_right = Inches(0.15)
        tf_c.margin_top = Inches(0.1)
        ph = tf_c.paragraphs[0]
        ph.alignment = PP_ALIGN.CENTER
        rh = ph.add_run()
        rh.text = title
        _set_font(rh, Pt(14), bold=True, color=accent)
        pd = tf_c.add_paragraph()
        pd.space_before = Pt(6)
        pd.alignment = PP_ALIGN.CENTER
        rd = pd.add_run()
        rd.text = desc
        _set_font(rd, Pt(12), color=CHARCOAL)

    _add_notes(slide, (
        "[~100 seconds]\n\n"
        "Let me walk you through the methodology. We parse digital sheet "
        "music -- MusicXML files -- using a Python toolkit called music21, "
        "and for each vocal line we build a tessituragram: every pitch "
        "maps to its total duration in quarter-note beats.\n\n"
        "The pipeline has five steps. The singer provides three inputs: "
        "their comfortable vocal range, favorite pitches, and pitches to "
        "avoid. We filter out any song whose written range goes beyond "
        "what the singer specified. Then we build an ideal pitch profile "
        "from those inputs. We create a list of numbers -- one per pitch "
        "in the singer's range. Every pitch starts with a small base "
        "weight so it is not ignored entirely. Favorite pitches get a "
        "large boost on top of that, and avoided pitches are dropped to "
        "zero. The result is a profile that peaks at the singer's "
        "preferred pitches and has nothing where they want to avoid.\n\n"
        "Each song has its own list built the same way -- but from the "
        "actual score, showing how much singing time falls on each pitch. "
        "Cosine similarity then compares the pattern of these two lists. "
        "If both concentrate time on the same pitches in the same "
        "proportions, the score is high -- close to 1.0. It focuses on "
        "shape rather than total duration. The avoid penalty is the "
        "proportion of the song's total singing duration that falls on "
        "notes the singer wants to avoid. Alpha controls the trade-off "
        "-- at 0.5, we split the weight evenly. Finally, we return a "
        "ranked list from best match to worst."
    ))


# ---------------------------------------------------------------------------
# Slide 5 -- Data & Evaluation Design  (merges old 8 + 9)
# ---------------------------------------------------------------------------

def slide_05_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "Data & Evaluation Design")
    _add_gold_accent_line(slide, Inches(1.18))

    gap = Inches(0.35)
    card_w = (Inches(12.2) - gap) / 2
    left_x = Inches(0.55)
    right_x = left_x + card_w + gap
    card_top = Inches(1.38)
    card_h = Inches(4.75)

    _two_col_card(
        slide, left_x, card_top, card_w, card_h,
        TEAL, LIGHT_GOLD,
        "OpenScore Lieder Corpus (CC0)",
        BURGUNDY,
        [
            "Exp 1: 101 vocal lines (one per composition).",
            "Exp 2: 1,655 lines from 1,419 compositions (~16\u00d7 larger).",
            "Multi-part works: each vocal line is its own item.",
        ],
        CHARCOAL,
    )

    _two_col_card(
        slide, right_x, card_top, card_w, card_h,
        SOFT_GRAY, LIGHT_GOLD,
        "Synthetic Self-Retrieval",
        BURGUNDY,
        [
            "Pick a vocal line from the library.",
            "Build a singer profile from that line\u2019s own tessituragram.",
            "Rank all candidates \u2014 does the same line come back first?",
        ],
        CHARCOAL,
    )

    # Baselines bar at bottom
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.65),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BURGUNDY
    bar.line.fill.background()
    tf_bar = bar.text_frame
    tf_bar.word_wrap = True
    tf_bar.margin_top = Inches(0.1)
    tf_bar.margin_left = Inches(0.3)
    pb = tf_bar.paragraphs[0]
    pb.alignment = PP_ALIGN.CENTER
    rb = pb.add_run()
    rb.text = "Three models:  Full (\u03B1 = 0.5)   |   Cosine-only (\u03B1 = 0)   |   Null (random)"
    _set_font(rb, Pt(16), bold=True, color=WHITE)

    _add_notes(slide, (
        "[~80 seconds]\n\n"
        "We used the OpenScore Lieder Corpus -- a freely available, "
        "openly licensed collection of art songs. Conveniently, all of "
        "the composers have been dead long enough that copyright is not "
        "a concern. We ran two experiments. The first used a compact "
        "library of 101 vocal lines, one per composition. The second "
        "used a much larger expanded library: 1,655 vocal lines drawn "
        "from 1,419 compositions -- about 16 times larger. Some "
        "compositions have multiple voice parts, like duets, so each "
        "vocal line is treated as its own item.\n\n"
        "We did not have human judges rating songs for us -- that is "
        "future work. Instead, we used a rigorous method called synthetic "
        "self-retrieval. We pick a vocal line from the library and build "
        "a singer profile directly from that line's own tessituragram: "
        "the range becomes the singer's range, the four pitches with the "
        "most singing time become the favorites, and the two pitches with "
        "the least become the avoids. Then we ask the system to rank all "
        "the remaining candidates and see where the original line ends "
        "up. If the system is working, it should rank that line very "
        "highly -- ideally first. We compared three models: the full "
        "model with the avoid penalty, cosine-only without the penalty, "
        "and a null baseline that just filters by range and then ranks "
        "randomly."
    ))


# ---------------------------------------------------------------------------
# Slide 6 -- Main Result: RQ1  (reuses old slide 10)
# ---------------------------------------------------------------------------

def slide_06_rq1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "Results: Self-Retrieval Accuracy (RQ1)")
    _add_gold_accent_line(slide, Inches(1.18))

    img = str(REPO / "paper_draft" / "figures" / "rq1_oracle_hr1_mrr.png")
    _add_image_centered(slide, img, top=Inches(1.45),
                        max_w=Inches(11.5), max_h=Inches(4.2))

    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(5.8), Inches(11.8), Inches(1.3),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = BURGUNDY
    shp.line.fill.background()
    # Clear bar text and use an aligned mini-table for direct comparison.
    shp.text_frame.clear()
    tbl_shape = slide.shapes.add_table(
        3, 4, Inches(1.15), Inches(5.93), Inches(10.95), Inches(1.02),
    )
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(4.35)
    tbl.columns[1].width = Inches(2.2)
    tbl.columns[2].width = Inches(2.2)
    tbl.columns[3].width = Inches(2.2)

    # Header row
    _format_table_cell(tbl.cell(0, 0), "Dataset", Pt(12), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(0, 1), "HR@1", Pt(12), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(0, 2), "HR@3", Pt(12), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(0, 3), "HR@5", Pt(12), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)

    # Data rows
    _format_table_cell(tbl.cell(1, 0), "Compact (101 lines)", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.LEFT, bg=BURGUNDY)
    _format_table_cell(tbl.cell(1, 1), "76%", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(1, 2), "98%", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(1, 3), "100%", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)

    _format_table_cell(tbl.cell(2, 0), "Expanded (1,655 lines)", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.LEFT, bg=BURGUNDY)
    _format_table_cell(tbl.cell(2, 1), "55%", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(2, 2), "80%", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(2, 3), "86%", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)

    _add_notes(slide, (
        "[~60 seconds]\n\n"
        "This figure shows the self-retrieval results. The metric on the "
        "left of each panel is Hit Rate at 1 -- how often the target song "
        "is ranked first -- and Mean Reciprocal Rank on the right, which "
        "captures how high it ranks on average. In Experiment 1, the "
        "compact 101-line library, hit rates are 76 percent at top 1, "
        "98 percent at top 3, and 100 percent at top 5. In Experiment 2 "
        "with 1,655 lines, hit rates are 55 percent at top 1, 80 percent "
        "at top 3, and 86 percent at top 5. These two experiments use "
        "different protocols and "
        "different candidate pools, so the drop from 76 to 55 percent "
        "is not purely a library-size effect. The key takeaway is clear: "
        "the system massively outperforms random ordering in both cases."
    ))


# ---------------------------------------------------------------------------
# Slide 7 -- Robustness: RQ2 + Alpha  (merges old 11 + 12)
# ---------------------------------------------------------------------------

def slide_07_robustness(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "Robustness: Stability & \u03B1 Sensitivity")
    _add_gold_accent_line(slide, Inches(1.18))

    # Left half: compact RQ2 table
    rows, cols = 5, 3
    tbl_left = Inches(0.35)
    tbl_top = Inches(1.35)
    tbl_w = Inches(6.3)
    tbl_h = Inches(2.9)
    tbl_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h)
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.55)
    tbl.columns[1].width = Inches(2.375)
    tbl.columns[2].width = Inches(2.375)

    c00 = tbl.cell(0, 0)
    c10 = tbl.cell(1, 0)
    c00.merge(c10)
    _format_table_cell(c00, "Model", Pt(10), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(0, 1), "Exp 1 (5 baselines)", Pt(10),
                       bold=True, color=WHITE, align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(0, 2), "Exp 2 (20 baselines)", Pt(10),
                       bold=True, color=WHITE, align=PP_ALIGN.CENTER, bg=BURGUNDY)
    sub = "Mean \u03C4 (95% CI)"
    _format_table_cell(tbl.cell(1, 1), sub, Pt(9), bold=True, color=SOFT_GRAY,
                       align=PP_ALIGN.CENTER, bg=LIGHT_GOLD)
    _format_table_cell(tbl.cell(1, 2), sub, Pt(9), bold=True, color=SOFT_GRAY,
                       align=PP_ALIGN.CENTER, bg=LIGHT_GOLD)

    data_rows = [
        ("Null", "\u22120.04 [\u22120.05, \u22120.02]", "0.00 [\u22120.002, 0.007]", LIGHT_GOLD),
        ("Cos-only (\u03B1=0)", "0.87 [0.84, 0.91]", "0.87 [0.86, 0.88]", None),
        ("Full (\u03B1=0.5)", "0.85 [0.81, 0.88]", "0.84 [0.82, 0.85]", LIGHT_GOLD),
    ]
    for ri, (m, v1, v2, stripe) in enumerate(data_rows, start=2):
        bg = stripe if stripe is not None else CREAM
        _format_table_cell(tbl.cell(ri, 0), m, Pt(10), bold=True, color=BURGUNDY,
                           align=PP_ALIGN.LEFT, bg=bg)
        _format_table_cell(tbl.cell(ri, 1), v1, Pt(10), bold=False, color=CHARCOAL,
                           align=PP_ALIGN.CENTER, bg=bg)
        _format_table_cell(tbl.cell(ri, 2), v2, Pt(10), bold=False, color=CHARCOAL,
                           align=PP_ALIGN.CENTER, bg=bg)

    cap_tb = _add_textbox(slide, Inches(0.35), Inches(4.35), Inches(6.3), Inches(0.5))
    cap_p = cap_tb.text_frame.paragraphs[0]
    cap_p.alignment = PP_ALIGN.LEFT
    cap_r = cap_p.add_run()
    cap_r.text = "Table 2. Mean Kendall\u2019s \u03C4; \u03C4 > 0.7 = strong agreement."
    _set_font(cap_r, Pt(11), italic=True, color=SOFT_GRAY)

    # Right: α vs mean τ only (HR@1/MRR panel omitted — covered on slide 6)
    img = str(REPO / "paper_draft" / "figures" / "alpha_sensitivity_hr1_mrr_tau.png")
    _add_image_tau_panel_from_alpha_figure(
        slide, img,
        left=Inches(6.75), top=Inches(1.35),
        max_w=Inches(6.35), max_h=Inches(3.55),
    )

    # Bottom summary bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.55), Inches(5.15), Inches(12.2), Inches(2.05),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BURGUNDY
    bar.line.fill.background()
    tf_bar = bar.text_frame
    tf_bar.word_wrap = True
    tf_bar.margin_top = Inches(0.12)
    tf_bar.margin_left = Inches(0.3)
    tf_bar.margin_right = Inches(0.3)

    summary_lines = [
        "Small preference changes \u2192 small ranking changes (\u03C4 = 0.84\u20130.85).",
        "Mean \u03C4 stays high across \u03B1 = 0\u20131 (\u03C4 \u2265 0.82 at \u03B1 = 1).",
        "\u03B1 = 0.5 balances avoid control and ranking stability.",
    ]
    for i, text in enumerate(summary_lines):
        p = tf_bar.paragraphs[0] if i == 0 else tf_bar.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = text
        _set_font(run, Pt(16), bold=(i == 2), color=WHITE)

    _add_notes(slide, (
        "[~90 seconds]\n\n"
        "So the system can find the right song -- but the next question "
        "is: is it consistent? If a singer tweaks just one note in their "
        "preferences, does the whole list of recommendations shuffle "
        "around, or does it mostly stay the same? That matters, because "
        "a system that gives you completely different answers every time "
        "you make a tiny change would not be very trustworthy.\n\n"
        "To test this, we made small, one-note changes to singer "
        "profiles and compared the recommendations before and after. "
        "We needed a measure that applies when the output is a ranked "
        "list of the same songs twice -- once before the edit and once "
        "after. Kendall's tau is a standard rank-correlation statistic "
        "for exactly that situation (Kendall, 1948): it looks at every "
        "pair of songs and checks whether they still appear in the same "
        "relative order in both lists. That is why it fits a stability "
        "question better than a top-one accuracy number alone -- Hit "
        "Rate at 1 only tells you whether the target song stayed first, "
        "whereas tau summarizes how much the entire ordering moved.\n\n"
        "Tau runs from 1.0 when the two lists are identical in order, "
        "down toward 0 when the rankings are unrelated, and can go "
        "negative if they tend to disagree. In our analysis we treat "
        "values above 0.7 as strong agreement between the two rankings. "
        "Our full model averaged about 0.84 -- well above that cutoff. "
        "The cosine-only variant was very similar, and the random "
        "baseline was near zero, as expected.\n\n"
        "The table on the left shows those consistency scores. On the "
        "right, the plot shows what happens when we change how much weight "
        "we give to the avoid penalty -- the alpha value from the "
        "formula earlier. We tried every value from 0 to 1. Self-"
        "retrieval accuracy was on the previous slide; here the curves "
        "show that ranking stability stays strong across the whole alpha "
        "range. We chose 0.5 as the default because it gives singers "
        "meaningful control over which notes to avoid without making the "
        "recommendations unstable."
    ))


# ---------------------------------------------------------------------------
# Slide 8 -- Takeaways & Limits  (merges old 14 + 15)
# ---------------------------------------------------------------------------

def slide_08_takeaways(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "Takeaways & Limitations")
    _add_gold_accent_line(slide, Inches(1.18))

    gap = Inches(0.35)
    card_w = (Inches(12.2) - gap) / 2
    left_x = Inches(0.55)
    right_x = left_x + card_w + gap
    card_top = Inches(1.38)
    card_h = Inches(4.45)

    _two_col_card(
        slide, left_x, card_top, card_w, card_h,
        TEAL, LIGHT_GOLD,
        "What does this mean?",
        BURGUNDY,
        [
            "In our tests, this gave a more consistent way to compare song fit.",
            "Looking at where notes are actually spent can reveal useful options.",
            "Simple, proven ranking tools can help in this new use case.",
        ],
        CHARCOAL,
    )

    _two_col_card(
        slide, right_x, card_top, card_w, card_h,
        GOLD, BURGUNDY,
        "Limitations & Future Work",
        GOLD,
        [
            "Synthetic profiles only \u2014 no human evaluation yet.",
            "One corpus (German & French art song); opera, theatre, pop untested.",
            "Pitch + duration only \u2014 dynamics, tempo, text omitted.",
            "Full model vs cosine-only was mixed in Exp 2; needs human testing.",
            "Next: real singers, diverse repertoire, richer features, interactive tool.",
        ],
        WHITE,
    )

    tb = _add_textbox(slide, Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = (
        "In offline tests, duration-weighted tessituragrams ranked songs by "
        "pitch-demand fit."
    )
    _set_font(r, Pt(20), italic=True, bold=True, color=BURGUNDY)

    _add_notes(slide, (
        "[~50 seconds]\n\n"
        "Let me step back and talk about what this means in practice. "
        "This is a proof-of-concept showing that data science can bring "
        "objectivity to a domain where decisions are traditionally "
        "subjective -- and where bad decisions have real health "
        "consequences. Looking at specific pitches gives a different kind of "
        "information than broad labels like Fach. It can surface pieces a teacher may not "
        "have considered that still look like a strong fit on these pitch "
        "patterns. For those of you in CS and data science, "
        "this also demonstrates that familiar tools -- cosine similarity, "
        "content-based filtering, offline evaluation -- can be useful in "
        "this vocal-repertoire setting and in a cold-start setting.\n\n"
        "I want to be upfront about limitations. These are synthetic "
        "profiles, not real singer preferences. We only tested on one "
        "corpus of German and French art song. And we only model pitch "
        "and duration -- dynamics, tempo, and text setting are not "
        "included. Also, in the larger experiment, the full model and the "
        "cosine-only model were often very close on the main retrieval "
        "metrics. For future work, we want to evaluate with real singers, "
        "expand to more diverse repertoire, add richer features, and "
        "ultimately build an interactive tool."
    ))


# ---------------------------------------------------------------------------
# Slide 9 -- Thank You
# ---------------------------------------------------------------------------

def slide_09_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BURGUNDY)

    for y in (Inches(1.5), Inches(5.8)):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(2), y, Inches(9.333), Inches(0.04),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = GOLD
        shp.line.fill.background()

    # Callback quote
    tb_q = _add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.333), Inches(1.2))
    tf_q = tb_q.text_frame
    tf_q.word_wrap = True
    pq = tf_q.paragraphs[0]
    pq.alignment = PP_ALIGN.CENTER
    rq = pq.add_run()
    rq.text = (
        "\u201cNext time a singer asks \u2018What should I sing?\u2019\n"
        "\u2014 data might have an answer.\u201d"
    )
    _set_font(rq, Pt(26), italic=True, color=LIGHT_GOLD)

    # Thank you
    tb_ty = _add_textbox(slide, Inches(1), Inches(3.6), Inches(11.333), Inches(1.0))
    pty = tb_ty.text_frame.paragraphs[0]
    pty.alignment = PP_ALIGN.CENTER
    rty = pty.add_run()
    rty.text = "Thank You"
    _set_font(rty, Pt(44), bold=True, color=GOLD)

    # Contact
    tb_c = _add_textbox(slide, Inches(1), Inches(4.8), Inches(11.333), Inches(0.9))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    contact = [
        "Madeline Johnson, Flint Million, Rajeev Bukralia",
        "Minnesota State University, Mankato",
    ]
    for i, line in enumerate(contact):
        p = tf_c.paragraphs[0] if i == 0 else tf_c.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        _set_font(r, Pt(20), color=LIGHT_GOLD)

    _add_notes(slide, (
        "[~15 seconds]\n\n"
        "To wrap up: duration-weighted tessituragrams can rank vocal "
        "repertoire by fit -- and this is just the beginning. Next time a "
        "singer asks 'What should I sing?', data might have an answer. "
        "Thank you very much for your time. I would be happy to take any "
        "questions."
    ))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_idea(prs)
    slide_04_method(prs)
    slide_05_data(prs)
    slide_06_rq1(prs)
    slide_07_robustness(prs)
    slide_08_takeaways(prs)
    slide_09_thanks(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
