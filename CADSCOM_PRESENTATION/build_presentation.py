"""
Build the CADSCOM 2026 conference presentation (.pptx).

Run from the repo root:
    python CADSCOM_PRESENTATION/build_presentation.py

Requires: python-pptx, Pillow
"""

from __future__ import annotations

import os
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = pathlib.Path(__file__).resolve().parent.parent
OUT = REPO / "CADSCOM_PRESENTATION" / "CADSCOM_2026_Presentation.pptx"

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
# Updated to CADSCOM 2026 blue-teal-coral scheme.
BURGUNDY   = RGBColor(0x0B, 0x1F, 0x3B)  # Deep Navy
CREAM      = RGBColor(0xFF, 0xFF, 0xFF)  # White
GOLD       = RGBColor(0xFF, 0x5D, 0x7A)  # Pink-toned coral accent
CHARCOAL   = RGBColor(0x1A, 0x1A, 0x1A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GOLD = RGBColor(0xF4, 0xF6, 0xF8)  # Light Gray
SOFT_GRAY  = RGBColor(0x2E, 0x4A, 0x7D)  # Slate Blue
TEAL       = RGBColor(0x3F, 0xA7, 0xA3)  # Soft Teal

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_burgundy_header(slide, height=Inches(1.15)):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = BURGUNDY
    shp.line.fill.background()
    return shp


def _add_gold_accent_line(slide, top, width=Inches(2.5)):
    left = Inches(0.75)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = GOLD
    shp.line.fill.background()


def _add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _set_font(run, size=Pt(22), bold=False, italic=False,
              color=CHARCOAL, name="Calibri"):
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = name


def _add_title_in_header(slide, text, top=Inches(0.25)):
    tb = _add_textbox(slide, Inches(0.75), top, Inches(11.5), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, Pt(32), bold=True, color=WHITE)
    p.alignment = PP_ALIGN.LEFT
    return tb


def _add_body_text(slide, lines: list[str], left=Inches(0.75),
                   top=Inches(1.5), width=Inches(11.8), height=Inches(5.5),
                   size=Pt(22), bullet=True, spacing=Pt(8)):
    tb = _add_textbox(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = ("  \u2022  " + line) if bullet else line
        _set_font(run, size, color=CHARCOAL)
        p.space_after = spacing
        p.alignment = PP_ALIGN.LEFT
    return tb


def _add_big_quote(slide, text, top=Inches(2.0)):
    tb = _add_textbox(slide, Inches(1.2), top, Inches(10.9), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    _set_font(run, Pt(34), italic=True, color=BURGUNDY)
    return tb


def _add_image_centered(slide, img_path: str, top=Inches(1.6),
                        max_w=Inches(11.0), max_h=Inches(5.3)):
    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    w = max_w
    h = int(w / aspect)
    if h > max_h:
        h = max_h
        w = int(h * aspect)
    left = int((SLIDE_W - w) / 2)
    slide.shapes.add_picture(img_path, left, top, w, h)


def _add_notes(slide, text: str):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def _format_table_cell(cell, text, size, bold=False, color=CHARCOAL,
                      align=PP_ALIGN.CENTER, bg=None):
    """Apply text and font styling to a table cell (python-pptx)."""
    cell.text = text
    tf = cell.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for p in tf.paragraphs:
        p.alignment = align
        for run in p.runs:
            _set_font(run, size, bold=bold, color=color)
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def _add_rq2_stability_table(slide):
    """Table 2 from paper / PDF: mean Kendall's τ (95% CI) by model and experiment."""
    rows, cols = 5, 3
    left, top = Inches(0.65), Inches(1.4)
    tbl_w, tbl_h = Inches(12.05), Inches(3.55)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, tbl_w, tbl_h)
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.35)
    tbl.columns[1].width = Inches(4.85)
    tbl.columns[2].width = Inches(4.85)

    c00 = tbl.cell(0, 0)
    c10 = tbl.cell(1, 0)
    c00.merge(c10)
    _format_table_cell(
        c00,
        "Model",
        Pt(13),
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        bg=BURGUNDY,
    )

    exp1_h = "Experiment 1 (5 baselines, 130 perturbations)"
    exp2_h = "Experiment 2 (20 baselines, 580 perturbations)"
    _format_table_cell(
        tbl.cell(0, 1), exp1_h, Pt(12), bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, bg=BURGUNDY,
    )
    _format_table_cell(
        tbl.cell(0, 2), exp2_h, Pt(12), bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, bg=BURGUNDY,
    )

    sub = "Mean \u03C4 (95% CI)"
    _format_table_cell(
        tbl.cell(1, 1), sub, Pt(12), bold=True, color=SOFT_GRAY,
        align=PP_ALIGN.CENTER, bg=LIGHT_GOLD,
    )
    _format_table_cell(
        tbl.cell(1, 2), sub, Pt(12), bold=True, color=SOFT_GRAY,
        align=PP_ALIGN.CENTER, bg=LIGHT_GOLD,
    )

    data_rows = [
        (
            "Null (random)",
            "\u22120.04 [\u22120.05, \u22120.02]",
            "0.00 [\u22120.002, 0.007]",
            LIGHT_GOLD,
        ),
        (
            "Cosine-only (\u03B1 = 0)",
            "0.87 [0.84, 0.91]",
            "0.87 [0.86, 0.88]",
            None,
        ),
        (
            "Full (\u03B1 = 0.5)",
            "0.85 [0.81, 0.88]",
            "0.84 [0.82, 0.85]",
            LIGHT_GOLD,
        ),
    ]
    for ri, (m, v1, v2, stripe) in enumerate(data_rows, start=2):
        bg = stripe if stripe is not None else CREAM
        _format_table_cell(
            tbl.cell(ri, 0), m, Pt(13), bold=True, color=BURGUNDY,
            align=PP_ALIGN.LEFT, bg=bg,
        )
        _format_table_cell(
            tbl.cell(ri, 1), v1, Pt(13), bold=False, color=CHARCOAL,
            align=PP_ALIGN.CENTER, bg=bg,
        )
        _format_table_cell(
            tbl.cell(ri, 2), v2, Pt(13), bold=False, color=CHARCOAL,
            align=PP_ALIGN.CENTER, bg=bg,
        )

    cap_tb = _add_textbox(
        slide, Inches(0.55), Inches(5.08), Inches(12.25), Inches(1.15),
    )
    cap_tf = cap_tb.text_frame
    cap_tf.word_wrap = True
    p1 = cap_tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = (
        "Table 2. Ranking stability: mean Kendall\u2019s \u03C4 per baseline (95% CI). "
        "\u03C4 = 1 means identical order; \u03C4 near 0 means unrelated rankings."
    )
    _set_font(r1, Pt(14), italic=True, color=SOFT_GRAY)
    p2 = cap_tf.add_paragraph()
    p2.space_before = Pt(8)
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "\u03C4 > 0.7 indicates strong agreement (Kendall, 1948)."
    _set_font(r2, Pt(15), bold=True, color=CHARCOAL)


def _standard_slide(prs, title, bullets, notes, **kw):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, title)
    _add_gold_accent_line(slide, Inches(1.18))
    if bullets:
        _add_body_text(slide, bullets, **kw)
    _add_notes(slide, notes)
    return slide

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BURGUNDY)

    # Gold accent lines
    for y in (Inches(1.8), Inches(5.6)):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(2), y, Inches(9.333), Inches(0.04),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = GOLD
        shp.line.fill.background()

    # Title
    tb = _add_textbox(slide, Inches(1), Inches(2.1), Inches(11.333), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Content-Based Vocal Repertoire Ranking Framework\nUsing Duration-Weighted Pitch Distributions"
    _set_font(run, Pt(34), bold=True, color=WHITE)

    # Authors
    tb2 = _add_textbox(slide, Inches(1), Inches(4.4), Inches(11.333), Inches(1.0))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Madeline Johnson, Flint Million, Rajeev Bukralia"
    _set_font(run2, Pt(24), color=LIGHT_GOLD)

    # Affiliation
    tb3 = _add_textbox(slide, Inches(1), Inches(5.05), Inches(11.333), Inches(0.5))
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Minnesota State University, Mankato"
    _set_font(run3, Pt(20), italic=True, color=LIGHT_GOLD)

    # Conference tag
    tb4 = _add_textbox(slide, Inches(1), Inches(5.9), Inches(11.333), Inches(0.5))
    p4 = tb4.text_frame.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run()
    run4.text = "CADSCOM 2026"
    _set_font(run4, Pt(20), bold=True, color=GOLD)

    _add_notes(slide, (
        "[~20 seconds]\n\n"
        "Good morning/afternoon, everyone. My name is Madeline Johnson, "
        "and I am currently pursuing a master's in data science here at "
        "Minnesota State University, Mankato. My undergraduate background "
        "is in vocal performance -- I hold a Bachelor of Music from "
        "Illinois Wesleyan University -- so this project sits right at "
        "the intersection of my two worlds. This work is co-authored with "
        "Flint Million and Rajeev Bukralia, and today I want to show you "
        "how we can use data from musical scores to help singers find "
        "repertoire that actually fits their voice."
    ))


def slide_02_hook(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)

    _add_big_quote(slide,
        "\u201cHave you ever wondered how a singer decides\n"
        "which songs are safe to sing?\u201d",
        top=Inches(1.8),
    )

    tb = _add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = [
        "The wrong piece is not just an inconvenience.",
        "It is an injury risk.",
        "What if data could help?",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        if i == 1:
            _set_font(run, Pt(26), bold=True, color=BURGUNDY)
        elif i == 2:
            _set_font(run, Pt(26), italic=True, color=GOLD)
        else:
            _set_font(run, Pt(24), color=CHARCOAL)
        p.space_after = Pt(12)

    # Subtle footnote citation
    fn = _add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4))
    fn_p = fn.text_frame.paragraphs[0]
    fn_p.alignment = PP_ALIGN.RIGHT
    fn_run = fn_p.add_run()
    fn_run.text = "Apfelbach, 2022; Phyland et al., 1999"
    _set_font(fn_run, Pt(12), italic=True, color=SOFT_GRAY)

    _add_notes(slide, (
        "[~45 seconds]\n\n"
        "As a trained singer, one of the first things I learned is that "
        "choosing the wrong piece is not just an inconvenience -- it can "
        "lead to real vocal injury. For example, a singer with a naturally "
        "lower voice will usually not feel as comfortable spending long "
        "periods in a high range as a singer with a naturally higher voice. "
        "Research shows that when a piece's demands do not match a singer's "
        "capabilities, the risk of strain goes up significantly (Apfelbach, "
        "2022; Phyland et al., 1999). In practice, singers and voice "
        "teachers rely on intuition, trial and error, or vocal "
        "classification systems that are not consistently defined. That is "
        "what motivated this research: using data to make the process of "
        "choosing repertoire more objective and safer."
    ))


def slide_03_current_approach(prs):
    """Slide 3: Range & Fach -- why current search filters fall short."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "The Current Approach")
    _add_gold_accent_line(slide, Inches(1.18))

    col_w = Inches(5.6)
    col_h = Inches(4.5)
    left_x = Inches(0.55)
    right_x = Inches(6.95)
    col_top = Inches(1.5)

    # --- Left column: Range ---
    box_l = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_x, col_top, col_w, col_h,
    )
    box_l.fill.solid()
    box_l.fill.fore_color.rgb = LIGHT_GOLD
    box_l.line.color.rgb = SOFT_GRAY
    box_l.line.width = Pt(1.2)
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.25)
    tf_l.margin_right = Inches(0.2)
    tf_l.margin_top = Inches(0.18)

    lh = tf_l.paragraphs[0]
    lh.alignment = PP_ALIGN.LEFT
    rh = lh.add_run()
    rh.text = "Filtering by Range"
    _set_font(rh, Pt(20), bold=True, color=SOFT_GRAY)

    left_items = [
        "Filter songs by minimum and\nmaximum pitch in the vocal line.",
        "Shows the extremes \u2014 not how\nmuch time the voice spends there.",
        "One high climax note looks the\nsame as sitting high the whole time.",
    ]
    for item in left_items:
        p = tf_l.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = "  \u2022  " + item
        _set_font(run, Pt(17), color=CHARCOAL)

    # --- Right column: Fach ---
    box_r = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_x, col_top, col_w, col_h,
    )
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = BURGUNDY
    box_r.line.fill.background()
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.25)
    tf_r.margin_right = Inches(0.2)
    tf_r.margin_top = Inches(0.18)

    rh2 = tf_r.paragraphs[0]
    rh2.alignment = PP_ALIGN.LEFT
    rr = rh2.add_run()
    rr.text = "Filtering by Fach"
    _set_font(rr, Pt(20), bold=True, color=GOLD)

    right_items = [
        "Categorizes voices by range, weight,\nand color (e.g. lyric soprano).",
        "Labels inconsistently defined across\nregions (Schloneger et al., 2024).",
        "Many voices fall between categories.",
        "Broad categories, not tailored to\nan individual singer.",
    ]
    for item in right_items:
        p = tf_r.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = "  \u2022  " + item
        _set_font(run, Pt(17), color=WHITE)

    # --- Bottom: bridging question ---
    tb = _add_textbox(slide, Inches(0.55), Inches(6.25), Inches(12.2), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Neither captures where the voice actually spends its time."
    _set_font(run, Pt(22), italic=True, color=BURGUNDY)

    _add_notes(slide, (
        "[~55 seconds]\n\n"
        "Today, singers mainly rely on two approaches. The first is "
        "filtering by range. Some databases let you filter by the notes "
        "in the vocal line -- you set a minimum and maximum pitch. That "
        "is useful, but it only tells you the extremes. A piece with one "
        "very high note at the climax looks the same as a piece that sits "
        "up high the entire time. Range tells you the boundaries, not "
        "where the voice actually lives.\n\n"
        "The second approach is Fach -- the system opera uses to categorize "
        "voices by range, weight, and color. Lyric soprano, dramatic mezzo, "
        "and so on. But Fach labels are not consistently defined across "
        "regions or pedagogical traditions (Schloneger et al., 2024), and "
        "many voices fall between categories -- which makes it hard to "
        "assign a single label. Fach gives you a broad category, not a "
        "profile tailored to a specific singer's unique voice. Neither "
        "approach captures where the voice actually spends its time -- "
        "and that distinction matters."
    ))


def slide_04_tessitura(prs):
    """Slide 4: Tessitura + Tessituragram (combined education slide)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "Tessitura & the Tessituragram")
    _add_gold_accent_line(slide, Inches(1.18))

    full_w = Inches(12.2)
    card_x = Inches(0.55)

    # ── Top card: Tessitura (compact definition) ──
    top_h = Inches(1.65)
    box_t = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, card_x, Inches(1.45), full_w, top_h,
    )
    box_t.fill.solid()
    box_t.fill.fore_color.rgb = LIGHT_GOLD
    box_t.line.color.rgb = SOFT_GRAY
    box_t.line.width = Pt(1.2)
    tf_t = box_t.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = Inches(0.35)
    tf_t.margin_right = Inches(0.3)
    tf_t.margin_top = Inches(0.18)

    ht = tf_t.paragraphs[0]
    ht.alignment = PP_ALIGN.LEFT
    rt = ht.add_run()
    rt.text = "Tessitura"
    _set_font(rt, Pt(24), bold=True, color=SOFT_GRAY)

    p_def = tf_t.add_paragraph()
    p_def.space_before = Pt(10)
    r_def = p_def.add_run()
    r_def.text = (
        "The pitches where the voice spends most of its time "
        "\u2014 not just the highest and lowest notes."
    )
    _set_font(r_def, Pt(20), color=CHARCOAL)

    # ── Bottom card: Tessituragram (larger, more detail) ──
    bot_top = Inches(3.35)
    bot_h = Inches(3.55)
    box_b = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, card_x, bot_top, full_w, bot_h,
    )
    box_b.fill.solid()
    box_b.fill.fore_color.rgb = BURGUNDY
    box_b.line.fill.background()
    tf_b = box_b.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = Inches(0.45)
    tf_b.margin_right = Inches(0.45)
    tf_b.margin_top = Inches(0.22)

    hb = tf_b.paragraphs[0]
    hb.alignment = PP_ALIGN.LEFT
    rb = hb.add_run()
    rb.text = "Tessituragram"
    _set_font(rb, Pt(24), bold=True, color=GOLD)

    bot_items = [
        "Thurmer (1988): histogram of singing time\nper pitch \u2014 a song's fingerprint.",
        "Duration-weighted: longer notes count more\n(Titze, 2008).",
        "Range and Fach don't quantify a song's\ninternal pitch-duration distribution.",
    ]
    for item in bot_items:
        p = tf_b.add_paragraph()
        p.space_before = Pt(18)
        run = p.add_run()
        run.text = "\u2022   " + item
        _set_font(run, Pt(19), color=WHITE)

    _add_notes(slide, (
        "[~65 seconds]\n\n"
        "There is a concept that does capture where the voice spends its "
        "time: tessitura. Range tells you the highest and lowest notes, "
        "but tessitura tells you where the voice actually lives -- the "
        "pitches where it spends most of its time. A piece might touch "
        "its highest note once at the climax, but if it sits in the "
        "middle of the range for most of its duration, that is where the "
        "real demand is. What matters most is not the extremes -- it is "
        "where you spend the bulk of your time.\n\n"
        "A researcher named Stefan Thurmer formalized a way to represent "
        "tessitura visually. In 1988 he introduced the tessituragram -- "
        "a histogram of singing time per pitch. Titze, in 2008, refined "
        "this with duration weighting, meaning longer notes count more "
        "because sustaining a pitch is more demanding than a passing "
        "tone. The result is a fingerprint of where the voice spends its "
        "time in that piece. This is exactly what range and Fach miss -- "
        "the internal distribution of vocal demand. That raises a natural "
        "follow-up: whether any existing tools already use tessituragrams "
        "this way."
    ))


def slide_05_gap(prs):
    """Slide 5: Existing tessitura tools & the gap they leave."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "The Gap")
    _add_gold_accent_line(slide, Inches(1.18))

    col_w = Inches(5.6)
    col_h = Inches(3.6)
    left_x = Inches(0.55)
    right_x = Inches(6.95)
    col_top = Inches(1.5)

    # --- Left: What exists ---
    box_l = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_x, col_top, col_w, col_h,
    )
    box_l.fill.solid()
    box_l.fill.fore_color.rgb = LIGHT_GOLD
    box_l.line.color.rgb = SOFT_GRAY
    box_l.line.width = Pt(1.2)
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.25)
    tf_l.margin_right = Inches(0.2)
    tf_l.margin_top = Inches(0.18)

    lh = tf_l.paragraphs[0]
    lh.alignment = PP_ALIGN.LEFT
    rh_run = lh.add_run()
    rh_run.text = "Existing Tessitura Tools"
    _set_font(rh_run, Pt(20), bold=True, color=SOFT_GRAY)

    left_items = [
        "Tessa (Apfelbach, 2022): extracts\ntessituragrams and summary statistics.",
        "Kassia Database: human-assessed\ntessitura per art song entry.",
        "Both evaluate a piece after it has\nalready been selected.",
    ]
    for item in left_items:
        p = tf_l.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = "  \u2022  " + item
        _set_font(run, Pt(17), color=CHARCOAL)

    # --- Right: What's missing ---
    box_r = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_x, col_top, col_w, col_h,
    )
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = BURGUNDY
    box_r.line.fill.background()
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.25)
    tf_r.margin_right = Inches(0.2)
    tf_r.margin_top = Inches(0.18)

    rh2 = tf_r.paragraphs[0]
    rh2.alignment = PP_ALIGN.LEFT
    rr = rh2.add_run()
    rr.text = "What\u2019s Missing"
    _set_font(rr, Pt(20), bold=True, color=GOLD)

    right_items = [
        "No tool recommends pieces whose\npitch distribution fits a singer.",
        "No pitch-level query: \u201cFind pieces\non THESE notes, avoid THOSE.\u201d",
        "217,000+ cataloged settings\n(LiederNet) \u2014 manual search can\u2019t scale.",
    ]
    for item in right_items:
        p = tf_r.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = "  \u2022  " + item
        _set_font(run, Pt(17), color=WHITE)

    # --- Bottom bridging line ---
    tb = _add_textbox(slide, Inches(0.55), Inches(5.4), Inches(12.2), Inches(1.5))
    tf_b = tb.text_frame
    tf_b.word_wrap = True
    p1 = tf_b.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "Tessituragrams capture vocal demand \u2014"
    _set_font(r1, Pt(21), italic=True, color=CHARCOAL)
    p2 = tf_b.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "but no tool uses them to recommend new repertoire."
    _set_font(r2, Pt(21), bold=True, italic=True, color=BURGUNDY)

    _add_notes(slide, (
        "[~55 seconds]\n\n"
        "There are tools that work with tessituragrams, but they are "
        "purely analytic. Tessa, developed by Apfelbach in 2022, extracts "
        "a tessituragram from a digital score and produces summary "
        "statistics. The Kassia Database -- a wonderful resource for art "
        "song by women composers -- displays a human-assessed tessitura "
        "for each entry. Both of these help you evaluate a piece after "
        "you have already found it. But neither one lets you use a "
        "tessituragram as a query -- to search for new pieces whose "
        "pitch distribution is likely to fit.\n\n"
        "The scale of the problem makes manual discovery impractical. "
        "The LiederNet Archive alone catalogs over 217,000 art-song "
        "settings. No voice teacher, no matter how experienced, can know "
        "all of that literature. Tessituragrams carry detailed information "
        "about a song's vocal demand -- but to our knowledge, no existing "
        "system uses them on the query side to recommend repertoire. "
        "That is the gap we wanted to explore."
    ))


def slide_06_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "The Pipeline")
    _add_gold_accent_line(slide, Inches(1.18))

    # Data-source line above the flowchart
    src_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(1.35), Inches(11.8), Inches(0.55),
    )
    src_box.fill.solid()
    src_box.fill.fore_color.rgb = LIGHT_GOLD
    src_box.line.color.rgb = TEAL
    src_box.line.width = Pt(1.2)
    tf_src = src_box.text_frame
    tf_src.word_wrap = True
    tf_src.margin_left = Inches(0.25)
    tf_src.margin_top = Inches(0.08)
    p_src = tf_src.paragraphs[0]
    p_src.alignment = PP_ALIGN.CENTER
    r_src = p_src.add_run()
    r_src.text = (
        "MusicXML scores parsed with music21: pitch \u2192 duration in quarter-note beats."
    )
    _set_font(r_src, Pt(17), italic=True, color=CHARCOAL)

    steps = [
        ("Singer\nPreferences", "Range, favorites,\nand avoids"),
        ("Range\nFilter", "Drop songs\noutside range"),
        ("Ideal Profile\nConstruction", "Target profile\nfrom preferences"),
        ("Cosine\nSimilarity\nScoring", "Similarity minus\navoid penalty"),
        ("Ranked\nSong List", "Best match\nto worst"),
    ]

    box_w = Inches(2.0)
    box_h = Inches(2.0)
    gap = Inches(0.45)
    total_w = len(steps) * box_w + (len(steps) - 1) * gap
    start_left = int((SLIDE_W - total_w) / 2)
    top_box = Inches(2.15)
    top_desc = Inches(4.35)

    for i, (label, desc) in enumerate(steps):
        left = start_left + i * (box_w + gap)

        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top_box, box_w, box_h,
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = BURGUNDY
        shp.line.fill.background()
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        _set_font(run, Pt(18), bold=True, color=WHITE)
        shp.text_frame.paragraphs[0].space_before = Pt(20)

        tb = _add_textbox(slide, left, top_desc, box_w, Inches(1.6))
        tf2 = tb.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        _set_font(run2, Pt(16), color=CHARCOAL)

        if i < len(steps) - 1:
            arrow_left = left + box_w
            arrow_top = top_box + box_h / 2 - Inches(0.15)
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top,
                gap, Inches(0.3),
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = TEAL
            arr.line.fill.background()

    # Bottom summary
    tb = _add_textbox(slide, Inches(1), Inches(6.2), Inches(11.333), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Three inputs: range, favorites, avoids. We do the rest."
    _set_font(run, Pt(22), italic=True, color=BURGUNDY)

    _add_notes(slide, (
        "[~50 seconds]\n\n"
        "Let me walk you through the pipeline. First, a quick note on "
        "the data. We parse digital sheet music -- MusicXML files -- "
        "using a Python toolkit called music21, and for each vocal line "
        "we build a tessituragram: every pitch maps to its total duration "
        "in quarter-note beats. That is the input to the system.\n\n"
        "The pipeline has five steps. The singer provides three inputs: "
        "their comfortable vocal range, favorite pitches, and pitches to "
        "avoid. We filter out any song whose written range goes beyond "
        "what the singer specified. Then we build an ideal pitch profile "
        "from those inputs -- the next slide explains exactly how that "
        "profile is formed. We score every "
        "remaining song by how closely its tessituragram matches that "
        "ideal, minus a penalty for time on avoided pitches. Finally, "
        "we return a ranked list from best match to worst."
    ))


def slide_07_formula(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "How Scoring Works")
    _add_gold_accent_line(slide, Inches(1.18))

    # ── Formula banner ──
    formula_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(1.45), Inches(11.8), Inches(1.0),
    )
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = BURGUNDY
    formula_box.line.fill.background()
    tf_f = formula_box.text_frame
    tf_f.word_wrap = True
    tf_f.margin_top = Inches(0.12)
    pf = tf_f.paragraphs[0]
    pf.alignment = PP_ALIGN.CENTER
    rf = pf.add_run()
    rf.text = "final_score  =  cosine_similarity(song, ideal)  \u2212  \u03B1  \u00D7  avoid_penalty"
    _set_font(rf, Pt(24), bold=True, color=WHITE, name="Consolas")

    # ── Three explanation cards ──
    card_w = Inches(3.7)
    card_h = Inches(2.8)
    card_top = Inches(2.7)
    card_gap = Inches(0.35)
    total_cards = 3 * card_w + 2 * card_gap
    card_start = int((SLIDE_W - total_cards) / 2)

    card_data = [
        (
            "Cosine Similarity",
            TEAL,
            "Same pitches, same proportions?\n\n"
            "1.0 = perfect match",
        ),
        (
            "Avoid Penalty",
            GOLD,
            "Share of singing time on\n"
            "pitches the singer avoids.\n\n"
            "Higher = worse fit",
        ),
        (
            "\u03B1  (Alpha = 0.5)",
            SOFT_GRAY,
            "Balances rewarding good\n"
            "pitches vs. penalizing\n"
            "avoided ones.\n\n"
            "0.5 = even split",
        ),
    ]

    for i, (label, accent, desc) in enumerate(card_data):
        left = card_start + i * (card_w + card_gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, card_w, card_h,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GOLD
        card.line.color.rgb = accent
        card.line.width = Pt(2.0)
        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.2)
        tf_c.margin_right = Inches(0.2)
        tf_c.margin_top = Inches(0.15)

        ph = tf_c.paragraphs[0]
        ph.alignment = PP_ALIGN.CENTER
        rh = ph.add_run()
        rh.text = label
        _set_font(rh, Pt(18), bold=True, color=accent)

        pd = tf_c.add_paragraph()
        pd.space_before = Pt(10)
        pd.alignment = PP_ALIGN.CENTER
        rd = pd.add_run()
        rd.text = desc
        _set_font(rd, Pt(15), color=CHARCOAL)

    _add_notes(slide, (
        "[~70 seconds]\n\n"
        "Here is how the scoring works. We create a list of numbers -- "
        "one per pitch in the singer's range. Every pitch starts with a "
        "small base weight so it is not ignored entirely. Favorite pitches "
        "get a large boost on top of that, and avoided pitches are dropped "
        "to zero. The result is a profile that peaks at the singer's "
        "preferred pitches and has nothing where they want to avoid.\n\n"
        "Each song has its own list built the same way -- one number per "
        "pitch -- but from the actual score, showing how much singing "
        "time falls on each pitch. Cosine similarity then compares the "
        "pattern of these two lists. If both concentrate time on the same "
        "pitches in the same proportions, the score is high -- close to "
        "1.0. It focuses on shape rather than total duration, so a short "
        "song and a long song that distribute time across pitches the "
        "same way score equally well. It is a standard tool for this kind "
        "of comparison in information retrieval, which is why we chose "
        "it.\n\n"
        "The avoid penalty is the proportion of the song's total singing "
        "duration that falls on notes the singer wants to avoid. Alpha "
        "controls the trade-off between rewarding good-fit pitches and "
        "penalizing avoided ones -- at 0.5, we split the weight evenly. "
        "The scoring function is straightforward in structure. To find "
        "out whether it actually produces good rankings, we ran two "
        "experiments on real musical data."
    ))


def slide_08_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "The Dataset")
    _add_gold_accent_line(slide, Inches(1.18))

    # ── Corpus banner ──
    corpus_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(1.45), Inches(11.8), Inches(0.85),
    )
    corpus_box.fill.solid()
    corpus_box.fill.fore_color.rgb = LIGHT_GOLD
    corpus_box.line.color.rgb = TEAL
    corpus_box.line.width = Pt(1.2)
    tf_c = corpus_box.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = Inches(0.3)
    tf_c.margin_top = Inches(0.1)
    pc = tf_c.paragraphs[0]
    pc.alignment = PP_ALIGN.CENTER
    rc = pc.add_run()
    rc.text = "OpenScore Lieder Corpus  \u2014  CC0-licensed art-song scores"
    _set_font(rc, Pt(19), italic=True, color=CHARCOAL)

    # ── Two experiment cards ──
    card_w = Inches(5.6)
    card_h = Inches(3.6)
    left_x = Inches(0.55)
    right_x = Inches(6.95)
    card_top = Inches(2.6)

    # --- Experiment 1 card ---
    box_l = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_x, card_top, card_w, card_h,
    )
    box_l.fill.solid()
    box_l.fill.fore_color.rgb = LIGHT_GOLD
    box_l.line.color.rgb = SOFT_GRAY
    box_l.line.width = Pt(1.2)
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.3)
    tf_l.margin_right = Inches(0.25)
    tf_l.margin_top = Inches(0.2)

    lh = tf_l.paragraphs[0]
    lh.alignment = PP_ALIGN.LEFT
    rl = lh.add_run()
    rl.text = "Experiment 1  \u2014  Compact Library"
    _set_font(rl, Pt(20), bold=True, color=SOFT_GRAY)

    exp1_items = [
        "101 vocal lines",
        "One per composition",
        "Initial validation baseline",
    ]
    for item in exp1_items:
        p = tf_l.add_paragraph()
        p.space_before = Pt(12)
        run = p.add_run()
        run.text = "  \u2022  " + item
        _set_font(run, Pt(18), color=CHARCOAL)

    # --- Experiment 2 card ---
    box_r = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_x, card_top, card_w, card_h,
    )
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = BURGUNDY
    box_r.line.fill.background()
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.3)
    tf_r.margin_right = Inches(0.25)
    tf_r.margin_top = Inches(0.2)

    rh2 = tf_r.paragraphs[0]
    rh2.alignment = PP_ALIGN.LEFT
    rr = rh2.add_run()
    rr.text = "Experiment 2  \u2014  Expanded Library"
    _set_font(rr, Pt(20), bold=True, color=GOLD)

    exp2_items = [
        "1,655 lines from 1,419 compositions",
        "~16\u00d7 larger",
        "342 multi-part, 1,313 single-voice",
    ]
    for item in exp2_items:
        p = tf_r.add_paragraph()
        p.space_before = Pt(12)
        run = p.add_run()
        run.text = "  \u2022  " + item
        _set_font(run, Pt(18), color=WHITE)

    # ── Bottom note ──
    tb = _add_textbox(slide, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.6))
    pb = tb.text_frame.paragraphs[0]
    pb.alignment = PP_ALIGN.CENTER
    rb = pb.add_run()
    rb.text = "Multi-part works (e.g. duets): each vocal line is its own item."
    _set_font(rb, Pt(17), italic=True, color=SOFT_GRAY)

    _add_notes(slide, (
        "[~35 seconds]\n\n"
        "We used the OpenScore Lieder Corpus -- a freely available, "
        "openly licensed collection of art songs. Conveniently, all of "
        "the composers have been dead long enough that copyright is not "
        "a concern. We ran two experiments. "
        "The first used a compact library of 101 vocal lines, one per "
        "composition. The second used a much larger expanded library: "
        "1,655 vocal lines drawn from 1,419 compositions -- about 16 "
        "times larger. Some compositions have multiple voice parts, like "
        "duets, so each vocal line is treated as its own item."
    ))


def slide_09_testing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "How We Tested It: Synthetic Self-Retrieval")
    _add_gold_accent_line(slide, Inches(1.18))

    # ── Two side-by-side cards (less visual noise than 3+3 mini-cards) ──
    gap = Inches(0.35)
    card_w = (Inches(12.2) - gap) / 2
    left_x = Inches(0.55)
    right_x = left_x + card_w + gap
    card_top = Inches(1.38)
    card_h = Inches(5.55)

    def _fill_two_column_card(box, border_rgb, title, lines, title_size=Pt(20)):
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GOLD
        box.line.color.rgb = border_rgb
        box.line.width = Pt(2.2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.28)
        tf.margin_right = Inches(0.26)
        tf.margin_top = Inches(0.22)

        ph = tf.paragraphs[0]
        ph.alignment = PP_ALIGN.LEFT
        rh = ph.add_run()
        rh.text = title
        _set_font(rh, title_size, bold=True, color=BURGUNDY)

        for line in lines:
            p = tf.add_paragraph()
            p.space_before = Pt(14)
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = "  \u2022  " + line
            _set_font(r, Pt(17), color=CHARCOAL)

    box_l = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_x, card_top, card_w, card_h,
    )
    _fill_two_column_card(
        box_l,
        TEAL,
        "Synthetic self-retrieval",
        [
            "Pick a vocal line from the library.",
            "Build a singer profile from that line\u2019s tessituragram "
            "(range, top-4 favorites, bottom-2 avoids).",
            "Rank all other candidates \u2014 does the same line rank first?",
        ],
    )

    box_r = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_x, card_top, card_w, card_h,
    )
    _fill_two_column_card(
        box_r,
        SOFT_GRAY,
        "Three models compared",
        [
            "Full:  cosine similarity + avoid penalty (\u03B1 = 0.5).",
            "Cosine-only:  similarity without the avoid penalty (\u03B1 = 0).",
            "Null:  range filter, then random ordering.",
        ],
    )

    _add_notes(slide, (
        "[~60 seconds]\n\n"
        "We did not have human judges rating songs for us -- that is "
        "future work. Instead, we used a rigorous method called synthetic "
        "self-retrieval. Here is how it works. We pick a vocal line from "
        "the library and build a singer profile directly from that line's "
        "own tessituragram: the range becomes the singer's range, the four "
        "pitches with the most singing time become the favorites, and the "
        "two pitches with the least become the avoids. Then we ask the "
        "system to rank all the remaining candidates and see where the "
        "original line ends up. If the system is working, it should rank "
        "that line very highly -- ideally first. We compared three models: "
        "the full model with the avoid penalty, cosine-only without the "
        "penalty, and a null baseline that just filters by range and then "
        "ranks randomly. Here is what we found."
    ))


def slide_10_rq1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "Results: Self-Retrieval Accuracy (RQ1)")
    _add_gold_accent_line(slide, Inches(1.18))

    img = str(REPO / "paper_draft" / "figures" / "rq1_oracle_hr1_mrr.png")
    _add_image_centered(slide, img, top=Inches(1.45), max_w=Inches(11.5), max_h=Inches(4.2))

    # Key numbers bar at bottom
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(5.8), Inches(11.8), Inches(1.3),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = BURGUNDY
    shp.line.fill.background()
    # Clear bar text and use an aligned mini-table for direct comparison.
    shp.text_frame.clear()
    tbl_shape = slide.shapes.add_table(
        3, 4, Inches(1.15), Inches(5.93), Inches(10.95), Inches(1.02),
    )
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(4.35)
    tbl.columns[1].width = Inches(2.2)
    tbl.columns[2].width = Inches(2.2)
    tbl.columns[3].width = Inches(2.2)

    # Header row
    _format_table_cell(tbl.cell(0, 0), "Dataset", Pt(12), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(0, 1), "HR@1", Pt(12), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(0, 2), "HR@3", Pt(12), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(0, 3), "HR@5", Pt(12), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)

    # Data rows
    _format_table_cell(tbl.cell(1, 0), "Compact (101 lines)", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.LEFT, bg=BURGUNDY)
    _format_table_cell(tbl.cell(1, 1), "76%", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(1, 2), "98%", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(1, 3), "100%", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)

    _format_table_cell(tbl.cell(2, 0), "Expanded (1,655 lines)", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.LEFT, bg=BURGUNDY)
    _format_table_cell(tbl.cell(2, 1), "55%", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(2, 2), "80%", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)
    _format_table_cell(tbl.cell(2, 3), "86%", Pt(14), bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, bg=BURGUNDY)

    _add_notes(slide, (
        "[~60 seconds]\n\n"
        "This figure shows the self-retrieval results. The metric on the "
        "left of each panel is Hit Rate at 1 -- how often the target song "
        "is ranked first -- and Mean Reciprocal Rank on the right, which "
        "captures how high it ranks on average. In Experiment 1, the "
        "compact 101-line library, hit rates are 76 percent at top 1, "
        "98 percent at top 3, and 100 percent at top 5. In Experiment 2 "
        "with 1,655 lines, hit rates are 55 percent at top 1, 80 percent "
        "at top 3, and 86 percent at top 5. These two experiments use "
        "different protocols and "
        "different candidate pools, so the drop from 76 to 55 percent "
        "is not purely a library-size effect. The key takeaway is clear: "
        "the system massively outperforms random ordering in both cases."
    ))


def slide_11_rq2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "Results: Ranking Stability (RQ2)")
    _add_gold_accent_line(slide, Inches(1.18))

    _add_rq2_stability_table(slide)

    _add_notes(slide, (
        "[~45 seconds]\n\n"
        "The next thing we wanted to know was whether the system is "
        "stable -- whether a small change to preferences causes the "
        "rankings to shift dramatically. We tested this with 580 one-note "
        "perturbations across 20 baseline profiles in the expanded "
        "library. We measured stability using Kendall's tau, a statistic "
        "that compares two ranked lists. A tau of 1.0 means identical "
        "rankings; 0.0 means completely unrelated. Anything above 0.7 is "
        "considered strong agreement. Our full model achieved a mean tau "
        "of 0.84. The cosine-only model was 0.87, with overlapping "
        "confidence intervals. The random baseline was essentially zero, "
        "as expected. When a singer tweaks their preferences slightly, "
        "the recommendations stay largely the same. The system is stable "
        "and trustworthy."
    ))


def slide_12_alpha(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "Sensitivity to the Avoid-Penalty Weight (\u03B1)")
    _add_gold_accent_line(slide, Inches(1.18))

    img = str(REPO / "paper_draft" / "figures" / "alpha_sensitivity_hr1_mrr_tau.png")
    _add_image_centered(slide, img, top=Inches(1.45), max_w=Inches(11.5), max_h=Inches(4.0))

    bullets = [
        "Performance flat across \u03B1 = 0\u20131.",
        "Stability holds: \u03C4 \u2265 0.82 even at \u03B1 = 1.0.",
        "\u03B1 = 0.5 balances avoid control and stability.",
    ]
    _add_body_text(slide, bullets, top=Inches(5.5), height=Inches(1.8), size=Pt(20))

    _add_notes(slide, (
        "[~45 seconds]\n\n"
        "We also tested how sensitive the results are to the choice of "
        "alpha -- the parameter that controls how heavily we penalize "
        "avoided notes. We swept alpha from 0 to 1 in both experiments. "
        "As you can see, self-retrieval performance -- Hit Rate at 1 and "
        "MRR -- is largely flat across the entire range. The system is "
        "not brittle to this choice. Stability does decrease slightly as "
        "alpha increases, which makes sense: a stronger avoid penalty "
        "creates more sensitivity to changes in avoid preferences. But "
        "even at alpha equals 1, tau stays above 0.82 -- well above the "
        "strong agreement threshold. We report alpha equals 0.5 as a "
        "balanced default: it gives singers meaningful control over their "
        "avoid preferences without making the rankings jittery."
    ))


def slide_13_rq3(prs):
    notes = (
        "[~30 seconds]\n\n"
        "As an engineering sanity check, we verified that the formula is "
        "implemented exactly as designed. The identity residual -- the "
        "difference between the computed score and what the formula "
        "predicts -- is exactly zero. An OLS regression recovers the "
        "exact coefficients: cosine weight equals 1.0, avoid weight "
        "equals negative 0.5, R-squared equals 1.0. All correlations go "
        "in the expected directions: higher cosine similarity predicts a "
        "higher final score, and higher avoid penalty predicts a lower "
        "score. No hidden bugs, no rounding surprises. The math checks "
        "out."
    )
    _standard_slide(prs,
        "Implementation Verification (RQ3)",
        [
            "Identity residual = 0  \u2014  zero numerical error.",
            "OLS recovers exact coefficients:\ncos = 1.0,  avoid = \u22120.5,  R\u00b2 = 1.0.",
            "Spearman \u03C1 in expected directions:\n"
            "(final, cos) +0.99  |  (final, avoid) \u22120.32  |  (cos, fav) +0.92",
            "The math checks out.",
        ],
        notes,
    )


def slide_14_so_what(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, CREAM)
    _add_burgundy_header(slide)
    _add_title_in_header(slide, "What Does This Mean?")
    _add_gold_accent_line(slide, Inches(1.18))

    bullets = [
        "Data science can make vocal health\ndecisions safer and more objective.",
        "Pitch-level queries give different\ninformation than broad labels like Fach.",
        "Familiar techniques, novel domain:\ncosine similarity, cold-start, offline eval.",
    ]
    _add_body_text(slide, bullets, top=Inches(1.5), height=Inches(3.0), size=Pt(22))

    _add_big_quote(slide,
        "\u201cHow might content-based recommendations\napply to creative domains like music?\u201d",
        top=Inches(4.8),
    )

    _add_notes(slide, (
        "[~50 seconds]\n\n"
        "Let me step back and talk about what this means in practice. "
        "This is a proof-of-concept showing that data science can bring "
        "objectivity to a domain where decisions are traditionally "
        "subjective -- and where bad decisions have real health "
        "consequences. Think about the practical impact for a voice "
        "teacher. Right now, they recommend pieces from their own training "
        "and experience. But a query based on specific pitches -- 'find me "
        "pieces that live on these notes and avoid those' -- gives a "
        "different kind of information than Fach labels. It could surface "
        "a piece the teacher has never encountered that still looks like a "
        "strong fit on these pitch patterns for their "
        "student's unique voice. For those of you in CS and data science, "
        "this also demonstrates that familiar tools -- cosine similarity, "
        "content-based filtering, offline evaluation -- can work in a "
        "completely new domain, and in a cold-start setting with no "
        "collaborative signal. I think there is a broader takeaway here "
        "for anyone working with content-based recommendation: these "
        "familiar techniques can open up entirely new domains, and I "
        "would welcome a conversation about where else they might apply."
    ))


def slide_15_limitations(prs):
    notes = (
        "[~30 seconds]\n\n"
        "I want to be upfront about limitations, because honesty "
        "strengthens research. First, these are synthetic profiles, not "
        "real singer preferences -- we have not done a human study yet. "
        "Second, we only tested on one corpus of German and French art "
        "song; opera, musical theatre, and popular song are untested. "
        "Third, we only model pitch and duration -- we do not account for "
        "dynamics, tempo, or text setting. For future work, we want to "
        "evaluate with real singers and their actual preferences, expand "
        "to more diverse repertoire, add richer musical features, and "
        "ultimately build an interactive tool that singers and voice "
        "teachers can use in practice."
    )
    _standard_slide(prs, "Limitations & Future Work", [
        "Synthetic profiles only \u2014 no human evaluation yet.",
        "One corpus (German & French art song);\nopera, theatre, pop untested.",
        "Pitch + duration only \u2014 dynamics,\ntempo, text omitted.",
        "Next: real singers, diverse repertoire,\nricher features, interactive tool.",
    ], notes)


def slide_16_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BURGUNDY)

    # Gold lines
    for y in (Inches(1.5), Inches(5.8)):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(2), y, Inches(9.333), Inches(0.04),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = GOLD
        shp.line.fill.background()

    # Summary
    tb = _add_textbox(slide, Inches(1), Inches(1.9), Inches(11.333), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = (
        "Duration-weighted tessituragrams can rank vocal repertoire\n"
        "by fit \u2014 and this is just the beginning."
    )
    _set_font(run, Pt(30), bold=True, color=WHITE)

    # Callback quote
    tb2 = _add_textbox(slide, Inches(1.5), Inches(3.7), Inches(10.333), Inches(1.0))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = (
        "\u201cNext time a singer asks \u2018What should I sing?\u2019\n"
        "\u2014 data might have an answer.\u201d"
    )
    _set_font(run2, Pt(24), italic=True, color=LIGHT_GOLD)

    # Thank you
    tb3 = _add_textbox(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.6))
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Thank You"
    _set_font(run3, Pt(36), bold=True, color=GOLD)

    # Contact
    tb4 = _add_textbox(slide, Inches(1), Inches(6.1), Inches(11.333), Inches(1.0))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    contact = [
        "Madeline Johnson, Flint Million, Rajeev Bukralia",
        "Minnesota State University, Mankato",
    ]
    for i, line in enumerate(contact):
        p4 = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        p4.alignment = PP_ALIGN.CENTER
        run4 = p4.add_run()
        run4.text = line
        _set_font(run4, Pt(18), color=LIGHT_GOLD)

    _add_notes(slide, (
        "[~15 seconds]\n\n"
        "To wrap up: duration-weighted tessituragrams can rank vocal "
        "repertoire by fit -- and this is just the beginning. Next time "
        "a singer asks 'What should I sing?', data might have an answer. "
        "Thank you very much for your time. I would be happy to take any "
        "questions."
    ))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_hook(prs)
    slide_03_current_approach(prs)
    slide_04_tessitura(prs)
    slide_05_gap(prs)
    slide_06_pipeline(prs)
    slide_07_formula(prs)
    slide_08_dataset(prs)
    slide_09_testing(prs)
    slide_10_rq1(prs)
    slide_11_rq2(prs)
    slide_12_alpha(prs)
    slide_13_rq3(prs)
    slide_14_so_what(prs)
    slide_15_limitations(prs)
    slide_16_closing(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
